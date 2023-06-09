from django.shortcuts import render, HttpResponse, redirect
import requests
from scApp import models
from pptx import Presentation
from django import forms
from django.core.validators import RegexValidator
from django.core.exceptions import ValidationError
from django.utils.safestring import mark_safe
from scApp.utils.pagination import Pagination
from scApp.utils.bootstrap import BootStrapModelForm
from scApp.utils.form import UserModelForm


# Create your views here.
def index(request):
    return HttpResponse("欢迎ssyi")


def tpl(request):
    name = 'ssyi'
    roles = ['管理员', 'CEO', '保安']
    user_info = {'name': 'ssyi', 'salary': 100000, 'role': 'CTO'}

    data_list = [
        {'name': 'ssyi01', 'salary': 200000, 'role': 'COO'},
        {'name': 'ssyi02', 'salary': 400000, 'role': 'CAO'},
        {'name': 'ssyi03', 'salary': 500000, 'role': 'CWO'},
    ]

    return render(request, 'tpl.html', {'n1': name, 'n2': roles, 'n3': user_info, 'n4': data_list})


def news(request):
    # 爬虫取新闻
    url = "http://www.chinaunicom.com.cn/api/article/NewsByIndex/2/2023/05/news"
    headers = {
        "Accept": "*/*",
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
                      "(HTML, like Gecko) Chrome/111.0.0.0 Safari/537.36"
    }
    res = requests.get(url, headers=headers)

    # data_list = res.json()
    # print(data_list)

    try:
        #  判断响应码是否为200
        if res.status_code == 200:
            data_list = res.json()

            #  判断数据结构是否正确
            if isinstance(data_list, list):
                # print(data_list)
                return render(request, 'news.html', {'new_list': data_list})
            else:
                print("返回数据错误")
        else:
            print("请求失败", res.status_code)
    except Exception as e:
        print(e)

    # return render(request, 'news.html', {'new_list': data_list})


def something(request):
    print(request.method)

    print(request.GET)

    # return HttpResponse('返回内容')

    return render(request, 'something.html', {'title': 'come on'})

    # return redirect("http://www.baidu.com")


def login(request):
    if request.method == 'GET':
        return render(request, 'login.html')
    else:
        # print(request.POST)
        username = request.POST.get("user")
        password = request.POST.get("pwd")

        if username == 'ssyi' and password == '000018':
            # return HttpResponse("登陆成功")
            return redirect("http://www.baidu.com")
        else:
            # return HttpResponse("用户名密码不正确!")
            return render(request, 'login.html', {'error_msg': '用户名或密码错误'})


# 操作数据库中数据
def orm(request):
    # models.Department.objects.create(title="市场部")
    # models.Department.objects.create(title="实施部")
    # models.UserInfo.objects.create(name='ssyi', password='000018', age='18')

    # 查询获取数据
    data_list = models.UserInfo.objects.filter(id=1)
    # data_list = models.UserInfo.objects.all()
    print(data_list)
    for obj in data_list:
        print(obj.id, obj.name)

    # 更新
    # models.UserInfo.objects.all().update(password=999)
    # models.UserInfo.objects.filter(id=1).update(password=999)
    return HttpResponse('成功')


def login_list(request):
    data_list = models.UserInfo.objects.all()
    return render(request, 'login_list.html', {'data_list': data_list})


def login_add(request):
    if request.method == 'GET':
        return render(request, 'login_add.html')

    user = request.POST.get("user")
    pwd = request.POST.get("pwd")
    age = request.POST.get("age")
    notes = request.POST.get("notes")

    # 验证用户输入
    if not user or not pwd or not age:
        return render(request, 'login_add.html', {'error': '必填项不可为空'})

    try:
        age = int(age)
    except ValueError:
        return render(request, 'login_add.html', {'error': '年龄必须为数字'})

    try:
        # 插入数据
        models.UserInfo.objects.create(name=user, password=pwd, age=age, notes=notes)
    except models.UserInfo.DoesNotExist:
        return render(request, 'login_add.html', {'error': '用户名已存在'})
    except Exception as e:
        # 出现异常，返回错误信息
        return render(request, 'login_add.html', {'error': '出现异常，请重试'})

    # 重定向到列表页面
    return redirect("/login/list/")


def login_del(request):
    nid = request.GET.get('nid')
    models.UserInfo.objects.filter(id=nid).delete()
    # return HttpResponse('删除成功')
    return redirect("/login/list/")


def ppt_view(request):
    prs = Presentation('scApp/path/ppt/file.pptx')

    # 获取PPT的幻灯片数量
    num_slides = len(prs.slides)

    # 获取第一张幻灯片，包括的文本框数量和内容
    slide = prs.slides[0]
    num_text_boxes = len(slide.shapes.placeholders)
    text_boxes_content = []
    for shape in slide.shapes.placeholders:
        if shape.has_text_frame:
            text_boxes_content.append(shape.text_frame.text)

    return render(request, 'ppt.html', {
        'num_slides': num_slides,
        'num_text_boxes': num_text_boxes,
        'text_boxes_content': text_boxes_content,
    })


def depart_list(request):
    queryset = models.Department.objects.all()
    return render(request, 'depart_list.html', {'queryset': queryset})


def depart_add(request):
    if request.method == 'GET':
        return render(request, 'depart_add.html')
    # 获取post提交内容
    name = request.POST.get('title')
    # 保存数据
    models.Department.objects.create(title=name)
    # 重定向到列表页面
    return redirect("/depart/list/")


def depart_del(request):
    nid = request.GET.get('nid')
    models.Department.objects.filter(id=nid).delete()
    # return HttpResponse('删除成功')
    return redirect("/depart/list/")


def depart_edit(request, nid):
    if request.method == 'GET':
        row_object = models.Department.objects.filter(id=nid).first()
        # 重定向到列表页面
        return render(request, 'depart_edit.html', {'row_object': row_object})

    # 获取post提交内容
    name = request.POST.get('title')
    models.Department.objects.filter(id=nid).update(title=name)
    return redirect("/depart/list/")


def user_list(request):
    # for i in range(30):
    #     models.UserInfo.objects.create(name='test002', password='123456', age='22', account='12',
    #                                    create_time='2023-05-22',
    #                                    gender='2', depart_id='2', role_id='2', )

    # 通过url传参数搜索查询功能实现
    data_dict = {}
    search_data = request.GET.get('query', '')
    if search_data:
        data_dict['name__contains'] = search_data

    # 搜索条件
    queryset = models.UserInfo.objects.filter(**data_dict).order_by('id')

    # 调用分页封装函数,实例化对象
    page_object = Pagination(request, queryset, page_size=9)

    # 定义一个字典
    context = {'queryset01': page_object.page_queryset, 'search_data': search_data,
               'page_string': page_object.html()}
    # 当前页
    # try:
    #     page = int(request.GET.get('page', 1))  # 分页
    # except  ValueError:
    #     page = 1
    # page_size = 10  # 每页显示数据
    # start = (page - 1) * page_size
    # end = page * page_size

    # 数据总条数
    # total_num = models.UserInfo.objects.filter(**data_dict).order_by('id').count()
    #
    # # 计算出总页码 divmod函数计算商和余数
    # total_page, div = divmod(total_num, page_object.page_size)
    # if div:
    #     total_page += 1

    # # 根据当前页计算出前后5页
    # plus = 5
    # if total_page <= 2 * plus + 1:
    #     # 当数据库数据少的时候，没有达到11页
    #     start_page = 1
    #     end_page = total_page
    # else:
    #     # 数据库数据比较多，判断当前页,<5
    #     if page <= plus:
    #         start_page = 1
    #         end_page = 2 * plus
    #     else:
    #         # 当前页大于5
    #         if (page + plus) > total_page:
    #             start_page = total_page - 2 * plus
    #             end_page = total_page
    #         else:
    #             start_page = page - plus
    #             end_page = page + plus
    #
    # # 页码
    # page_list = []
    # page_list.append('<li><a href="?page={}">首页</a></li>'.format(1))
    # # 上一页
    # if page > 1:
    #     prev = '<li><a href="?page={}">上一页</a></li>'.format(page - 1)
    # else:
    #     prev = '<li><a href="?page={}">上一页</a></li>'.format(1)
    # page_list.append(prev)
    #
    # # 分页页面
    # for i in range(start_page, end_page + 1):
    #     if i == page:
    #         ele = '<li class="active"><a href="?page={}">{}</a></li>'.format(i, i)
    #     else:
    #         ele = '<li><a href="?page={}">{}</a></li>'.format(i, i)
    #     page_list.append(ele)
    #
    # # 下一页
    # if page < total_page:
    #     prev = '<li><a href="?page={}">下一页</a></li>'.format(page + 1)
    # else:
    #     prev = '<li><a href="?page={}">下一页</a></li>'.format(total_page)
    # page_list.append(prev)
    #
    # page_list.append('<li><a href="?page={}">尾页</a></li>'.format(total_page))
    #
    # search_tz = '''
    #              <li>
    #                 <form style="float: left;margin-left: -1px" method="get">
    #                     <input name="page"
    #                            style="position: relative;float: left;display: inline-block;width: 80px;border-radius: 0;"
    #                            type="text" class="form-control" placeholder="页码">
    #                     <button style="border-radius: 0" class="btn btn-default" type="submit">跳转</button>
    #                     </span>
    #                 </form>
    #             </li>
    # '''
    # page_list.append(search_tz)
    #
    # # 导入mark_safe，把数据包裹成安全的传递给html
    # page_string = mark_safe(''.join(page_list))

    # queryset = models.UserInfo.objects.filter(name__contains=search_data)
    # print(queryset)

    # for obj in queryset:
    #     # obj.get_gender_display() 自动找定义的元组数据
    #     # obj.depart.title  # 获取关联表数据
    #     print(obj.create_time.strftime('%Y-%m-%d'), obj.get_gender_display(),
    #           obj.depart.title
    #           )
    # queryset = models.UserInfo.objects.filter(**data_dict).order_by('id')[page_object.start:page_object.end]

    return render(request, 'user_list.html', context)


def user_add(request):
    if request.method == 'GET':
        context = {
            'gender_choices': models.UserInfo.gender_choices,
            'depart_list': models.Department.objects.all(),
            'role_list': models.Role.objects.all(),
        }
        return render(request, 'user_add.html', context)

    # post获取页面value里数据
    user = request.POST.get('user')
    pwd = request.POST.get('pwd')
    age = request.POST.get('age')
    ac = request.POST.get('ac')
    ctime = request.POST.get('ctime')
    gd = request.POST.get('gd')
    dp = request.POST.get('dp')
    re = request.POST.get('re')

    models.UserInfo.objects.create(name=user, password=pwd, age=age, account=ac, create_time=ctime,
                                   gender=gd, depart_id=dp, role_id=re, )
    return redirect("/user/list/")


# 使用modelform生成页面input框,编辑class可以另外再写一个，需要用的进行实例化即可
# class UserModelForm(forms.ModelForm):

# 继承自定义的modelform样式
class UserModelForm(BootStrapModelForm):
    # 输入长度验证规则
    name = forms.CharField(min_length=5,
                           label='姓名',
                           widget=forms.TextInput(attrs={'class': "form-control"}))
    # name = forms.CharField(disabled=True,label='姓名')  # 字段不可修改，或者直接在下面fields里去掉这个字段名即可
    # 正则表达式输入验证(方式一)
    age = forms.CharField(label='年龄啊',
                          validators=[RegexValidator(r'^[0-9]+$', '年龄请输入数字')]
                          )

    class Meta:
        model = models.UserInfo
        fields = ['name', 'password', 'age', 'account', 'create_time', 'gender', 'depart', 'role']
        # fields = '__all__'  # 表示所有的字段
        # exclude = ['age']  # 排除某个字段

        # 插件里加样式
        # widgets = {
        #     'name': forms.TextInput(attrs={'class': "form-control"}),
        #     'password': forms.PasswordInput(attrs={'class': "form-control"}),
        # }

    # 重写方法，循环找所有插件，进行重写class样式
    # def __init__(self, *args, **kwargs):
    #     super(UserModelForm, self).__init__(*args, **kwargs)
    #
    #     for name, field in self.fields.items():
    #         # print(name, field)
    #         # 单独设置role角色不显示样式
    #         if name == 'role':
    #             continue
    #         field.widget.attrs = {'class': "form-control", 'placeholder': field.label}

    # 验证方式二(钩子方法)
    def clean_password(self):
        txt_pwd = self.cleaned_data['password']
        if len(txt_pwd) != 6:
            raise ValidationError("输入的位数错误")
        # 验证通过
        return txt_pwd

    # 判断姓名不允许重复(新建和编辑)
    def clean_name(self):
        # 排除当前编辑的一行ID exclude(self.instance.pk)

        txt_name = self.cleaned_data['name']
        exname = models.UserInfo.objects.exclude(id=self.instance.pk).filter(name=txt_name).exists()
        if exname:
            raise ValidationError('姓名不允许重复')
        return txt_name


def user_modelform_add(request):
    if request.method == 'GET':
        form = UserModelForm()
        return render(request, 'user_modelform_add.html', {'form': form})

    # POST进行提交后校验
    form = UserModelForm(data=request.POST)
    if form.is_valid():
        # print(form.cleaned_data)
        form.save()
        return redirect("/user/list/")
    else:
        return render(request, 'user_modelform_add.html', {'form': form})


def user_modelform_edit(request, nid):
    # 根据ID去数据库获取那一行数据（对象）
    row_object = models.UserInfo.objects.filter(id=nid).first()

    if request.method == 'GET':
        form1 = UserModelForm(instance=row_object)  # 实例化后获取默认值
        return render(request, 'user_modelform_edit.html', {'form': form1})  # 把数据传到html里
    # 更新这一行数据
    form = UserModelForm(data=request.POST, instance=row_object)
    if form.is_valid():
        # 如果不是界面需要输入的值保存，可以用
        # form.instance.字段名 = '某个值'
        form.save()
        return redirect("/user/list/")
    else:
        return render(request, 'user_modelform_add.html', {'form': form})


def user_modelform_del(request, nid):
    models.UserInfo.objects.filter(id=nid).delete()
    return redirect("/user/list/")
